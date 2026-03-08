"""
CIVITAS Pitch Deck Generator

Generates a 13-slide PowerPoint pitch deck targeting small real estate
law firms and title companies in Chicago.

Usage:
    pip install python-pptx
    python scripts/generate_pitch_deck.py
    # -> civitas_pitch_deck.pptx in project root
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# -- Design tokens (minimal white/grey palette) ------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFF, 0xFF, 0xFF)           # pure white slides
BG_ALT = RGBColor(0xF9, 0xFA, 0xFB)      # gray-50 for subtle contrast
CARD_BG = RGBColor(0xF9, 0xFA, 0xFB)     # gray-50 card fill
BORDER = RGBColor(0xE5, 0xE7, 0xEB)      # gray-200
BORDER_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)  # gray-100

CHARCOAL = RGBColor(0x11, 0x18, 0x27)    # near-black headings
TEXT_BODY = RGBColor(0x4B, 0x55, 0x63)    # gray-600
TEXT_SEC = RGBColor(0x9C, 0xA3, 0xAF)     # gray-400
TEXT_LIGHT = RGBColor(0xD1, 0xD5, 0xDB)   # gray-300

DARK = RGBColor(0x1F, 0x25, 0x37)         # gray-800 for hero sections
DARK_MID = RGBColor(0x37, 0x41, 0x51)     # gray-700

FONT = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Activity levels use a subtle grey scale
LEVEL_COLORS = {
    "QUIET": RGBColor(0xD1, 0xD5, 0xDB),     # gray-300
    "TYPICAL": RGBColor(0x9C, 0xA3, 0xAF),   # gray-400
    "ACTIVE": RGBColor(0x4B, 0x55, 0x63),     # gray-600
    "COMPLEX": RGBColor(0x1F, 0x25, 0x37),    # gray-800
}


# -- Helpers ------------------------------------------------------------------

def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rule(slide, top, left=0.8, width=11.7):
    """Thin horizontal rule / divider line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Pt(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER
    shape.line.fill.background()


def _add_footer(slide, page_num=None):
    # Divider above footer
    _add_rule(slide, 6.75, 0.6, 12.1)

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(6.9), Inches(2), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "CIVITAS"
    run.font.name = FONT
    run.font.size = Pt(9)
    run.font.color.rgb = TEXT_SEC
    run.font.bold = True
    p.space_before = Pt(0)
    p.space_after = Pt(0)

    if page_num is not None:
        txBox2 = slide.shapes.add_textbox(
            Inches(11.8), Inches(6.9), Inches(1), Inches(0.4)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        run2 = p2.add_run()
        run2.text = str(page_num)
        run2.font.name = FONT
        run2.font.size = Pt(9)
        run2.font.color.rgb = TEXT_SEC


def _add_title(slide, text, left=0.8, top=0.8, width=11.5, size=30):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(0.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = CHARCOAL
    run.font.bold = True
    return txBox


def _add_subtitle(slide, text, left=0.8, top=1.5, width=11.5, size=15):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(0.6)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = TEXT_SEC
    return txBox


def _add_body(slide, text, left=0.8, top=2.2, width=11.5, height=3.5, size=14):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = TEXT_BODY
    p.line_spacing = Pt(24)
    return txBox


def _add_card(slide, left, top, width, height, title, body,
              number=None, title_size=15, body_size=12):
    """Minimal card with light grey fill and subtle border."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    shape.adjustments[0] = 0.04

    y_offset = 0.2

    # Optional large number (for numbered cards)
    if number is not None:
        txNum = slide.shapes.add_textbox(
            Inches(left + 0.25), Inches(top + 0.15),
            Inches(0.5), Inches(0.4)
        )
        tf_n = txNum.text_frame
        p_n = tf_n.paragraphs[0]
        run_n = p_n.add_run()
        run_n.text = str(number)
        run_n.font.name = FONT
        run_n.font.size = Pt(22)
        run_n.font.color.rgb = TEXT_LIGHT
        run_n.font.bold = True
        y_offset = 0.55

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.25), Inches(top + y_offset),
        Inches(width - 0.4), Inches(0.4)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT
    run.font.size = Pt(title_size)
    run.font.color.rgb = CHARCOAL
    run.font.bold = True

    # Body
    body_top = y_offset + 0.35
    txBox2 = slide.shapes.add_textbox(
        Inches(left + 0.25), Inches(top + body_top),
        Inches(width - 0.4), Inches(height - body_top - 0.15)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = body
    run2.font.name = FONT
    run2.font.size = Pt(body_size)
    run2.font.color.rgb = TEXT_BODY
    p2.line_spacing = Pt(19)


def _add_level_pill(slide, left, top, level, width=1.6, height=0.42):
    color = LEVEL_COLORS.get(level, TEXT_SEC)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.4
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = level
    run.font.name = FONT
    run.font.size = Pt(13)
    run.font.color.rgb = WHITE
    run.font.bold = True


def _add_stat_card(slide, left, top, number, label):
    """Clean stat card with large number on grey background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(3.4), Inches(2.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    shape.adjustments[0] = 0.04

    # Number
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top + 0.35), Inches(3.4), Inches(0.8)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = number
    run.font.name = FONT
    run.font.size = Pt(38)
    run.font.color.rgb = CHARCOAL
    run.font.bold = True

    # Label
    txBox2 = slide.shapes.add_textbox(
        Inches(left + 0.3), Inches(top + 1.25), Inches(2.8), Inches(0.7)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.name = FONT
    run2.font.size = Pt(12)
    run2.font.color.rgb = TEXT_SEC


def _add_step_box(slide, left, top, number, title, desc):
    """Numbered step with grey circle."""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + 0.5), Inches(top), Inches(0.7), Inches(0.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.name = FONT
    run.font.size = Pt(22)
    run.font.color.rgb = WHITE
    run.font.bold = True

    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top + 0.85), Inches(1.7), Inches(0.4)
    )
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = title
    run2.font.name = FONT
    run2.font.size = Pt(14)
    run2.font.color.rgb = CHARCOAL
    run2.font.bold = True

    txBox2 = slide.shapes.add_textbox(
        Inches(left - 0.15), Inches(top + 1.25), Inches(2), Inches(0.8)
    )
    tf3 = txBox2.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = desc
    run3.font.name = FONT
    run3.font.size = Pt(11)
    run3.font.color.rgb = TEXT_SEC


def _add_arrow(slide, left, top):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(0.7), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER
    shape.line.fill.background()


def _add_bullet_list(slide, items, left=0.8, top=2.2, width=11.5, size=15):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(4)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(10)
        p.line_spacing = Pt(26)
        run = p.add_run()
        run.text = f"\u2022  {item}"
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = TEXT_BODY


# -- Slide builders -----------------------------------------------------------

def _slide_01_title(prs):
    """Title slide: dark charcoal hero, white below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)

    # Dark hero band
    hero = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(3.6)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = DARK
    hero.line.fill.background()

    # CIVITAS wordmark
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.8), Inches(11.5), Inches(1.2)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "C I V I T A S"
    run.font.name = FONT
    run.font.size = Pt(54)
    run.font.color.rgb = WHITE
    run.font.bold = True

    # Subtitle
    txBox2 = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.6)
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Municipal Intelligence"
    run2.font.name = FONT
    run2.font.size = Pt(20)
    run2.font.color.rgb = TEXT_LIGHT

    # Tagline below hero
    txBox3 = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.2), Inches(11.5), Inches(1)
    )
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Transaction-grade property intelligence\nin minutes, not days."
    run3.font.name = FONT
    run3.font.size = Pt(24)
    run3.font.color.rgb = CHARCOAL

    # Chicago v1 note
    txBox4 = slide.shapes.add_textbox(
        Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.5)
    )
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    run4 = p4.add_run()
    run4.text = "Chicago v1  \u2022  Built for real estate law firms & title companies"
    run4.font.name = FONT
    run4.font.size = Pt(13)
    run4.font.color.rgb = TEXT_SEC

    _add_footer(slide)


def _slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "The Problem")
    _add_subtitle(slide, "Manual municipal due diligence is slow, fragmented, and error-prone.")
    _add_rule(slide, 2.1)

    cards = [
        ("Scattered Data",
         "Law firms manually search 6+ city databases\u2014violations, permits, inspections, "
         "tax liens, 311 complaints\u2014each with different interfaces and formats."),
        ("No Unified Picture",
         "There is no single view connecting a property's violation history, permit status, "
         "tax standing, and complaint patterns. Critical context gets missed."),
        ("Missed Findings",
         "Without systematic analysis, aged violations, recurring complaints, and tax lien "
         "patterns fall through the cracks\u2014creating liability exposure."),
    ]
    for i, (title, body) in enumerate(cards):
        _add_card(slide, 0.8 + i * 4.0, 2.5, 3.6, 3.0, title, body, number=i + 1)

    _add_footer(slide, 2)


def _slide_03_cost(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "The Cost of Inaction")
    _add_subtitle(slide, "What happens when due diligence gaps go unaddressed.")
    _add_rule(slide, 2.1)

    _add_stat_card(slide, 0.5, 2.5, "4\u20136 hrs", "Average time per manual\nmunicipal records search")
    _add_stat_card(slide, 4.8, 2.5, "23%", "Of Chicago properties have\nat least one open violation")
    _add_stat_card(slide, 9.1, 2.5, "$50K+", "Potential liability from\nundiscovered liens & violations")

    _add_body(slide,
              "Every missed violation or undisclosed lien is a potential post-closing dispute. "
              "Title companies and law firms bear the reputational and financial risk.",
              top=5.2, size=13)

    _add_footer(slide, 3)


def _slide_04_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "The Solution: CIVITAS")
    _add_subtitle(slide, "Automated municipal intelligence for real estate transactions.")
    _add_rule(slide, 2.1)

    features = [
        "Aggregates 6 municipal & tax datasets into a unified property profile",
        "Applies 15 deterministic, auditable rules across 4 action categories",
        "Computes transparent, weighted activity scores with full citation",
        "Generates AI-narrated executive summaries grounded in structured data",
        "Delivers professional PDF reports in under 60 seconds",
    ]
    _add_bullet_list(slide, features, top=2.4, size=15)

    _add_footer(slide, 4)


def _slide_05_how_it_works(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "How It Works")
    _add_subtitle(slide, "From address to actionable intelligence in three steps.")
    _add_rule(slide, 2.1)

    _add_step_box(slide, 2.0, 2.8, 1, "Enter Address",
                  "Type a Chicago address or PIN.\n6-tier resolution matches\nto canonical records.")
    _add_arrow(slide, 4.2, 3.0)
    _add_step_box(slide, 5.3, 2.8, 2, "Automated Analysis",
                  "SQL rule engine evaluates\n15 rules across 6 datasets\ninstantly.")
    _add_arrow(slide, 7.5, 3.0)
    _add_step_box(slide, 8.6, 2.8, 3, "Property Report",
                  "Scored report with findings,\nsupporting records, AI\nnarrative, and PDF export.")

    _add_body(slide,
              "Every finding is deterministic and citable. Claude AI interprets structured results\u2014"
              "it never invents findings or computes scores.",
              top=5.3, size=12)

    _add_footer(slide, 5)


def _slide_06_activity_scoring(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "Activity Scoring Model")
    _add_subtitle(slide, "Weighted, additive, and fully transparent.")
    _add_rule(slide, 2.1)

    levels = [
        ("QUIET", "0 \u2013 24", "No active findings"),
        ("TYPICAL", "25 \u2013 49", "Minor compliance items"),
        ("ACTIVE", "50 \u2013 74", "Notable municipal activity"),
        ("COMPLEX", "75+", "Multiple findings requiring attention"),
    ]
    for i, (level, score_range, desc) in enumerate(levels):
        y = 2.5 + i * 0.85
        _add_level_pill(slide, 1.0, y, level)
        txBox = slide.shapes.add_textbox(Inches(3.0), Inches(y + 0.05), Inches(1.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = score_range
        run.font.name = FONT
        run.font.size = Pt(15)
        run.font.color.rgb = CHARCOAL
        run.font.bold = True
        txBox2 = slide.shapes.add_textbox(Inches(4.8), Inches(y + 0.05), Inches(4), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = desc
        run2.font.name = FONT
        run2.font.size = Pt(13)
        run2.font.color.rgb = TEXT_BODY

    _add_card(slide, 9.2, 2.3, 3.5, 3.5,
              "4 Action Groups",
              "Action Required \u2014 Tax & financial (25\u201340 pts)\n"
              "Review Recommended \u2014 Enforcement (20\u201335 pts)\n"
              "Worth Noting \u2014 Compliance (15\u201320 pts)\n"
              "Informational \u2014 Regulatory friction (10\u201315 pts)\n\n"
              "15 rules total \u2022 Configurable weights",
              body_size=11)

    _add_footer(slide, 6)


def _slide_07_report_deep_dive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "Report Deep-Dive")
    _add_subtitle(slide, "Everything a closing attorney needs in one document.")
    _add_rule(slide, 2.1)

    sections = [
        ("Executive Summary", "Activity score, level, triggered\nfindings, and AI-generated overview."),
        ("Findings by Action Group", "Action Required, Review Recommended,\nWorth Noting, Informational."),
        ("Supporting Records", "Violations, inspections, permits,\n311 requests, tax liens\u2014sortable tables."),
        ("AI Narrative", "Claude interprets findings in legally\ncautious, citation-backed language."),
        ("PDF Export", "Professional report ready to attach\nto closing files or send to clients."),
        ("Legal Disclaimer", "Clear statement that CIVITAS does\nnot replace formal title examination."),
    ]
    for i, (title, body) in enumerate(sections):
        col = i % 3
        row = i // 3
        _add_card(slide, 0.8 + col * 4.0, 2.4 + row * 2.2, 3.6, 1.8,
                  title, body, title_size=13, body_size=11)

    _add_footer(slide, 7)


def _slide_08_data_sources(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "Data Sources")
    _add_subtitle(slide, "Six authoritative Chicago datasets power every report.")
    _add_rule(slide, 2.1)

    datasets = [
        ("Building Violations", "Dept. of Buildings enforcement\nactions and compliance orders."),
        ("Food Inspections", "Health dept. inspection results\nand pass/fail/conditional outcomes."),
        ("Building Permits", "Active, completed, and delayed\npermit applications and status."),
        ("311 Service Requests", "13.4M+ citizen complaints\nincluding building and sanitation."),
        ("Tax Liens", "Cook County Clerk records\nof delinquent property tax liens."),
        ("Vacant Buildings", "Registered vacant/abandoned\nbuildings on the Chicago Data Portal."),
    ]
    for i, (title, body) in enumerate(datasets):
        col = i % 3
        row = i // 3
        _add_card(slide, 0.8 + col * 4.0, 2.4 + row * 2.2, 3.6, 1.8,
                  title, body, number=i + 1, title_size=13, body_size=11)

    _add_footer(slide, 8)


def _slide_09_portfolio(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "Portfolio Analysis")
    _add_subtitle(slide, "Analyze entire portfolios at once with batch CSV upload.")
    _add_rule(slide, 2.1)

    features = [
        ("CSV Upload", "Upload a list of addresses.\nEach property is analyzed\nagainst all 15 rules."),
        ("Live Progress", "Server-sent events stream\nreal-time status as each\nproperty is processed."),
        ("Summary Dashboard", "Average activity score, level\ndistribution chart, and\nper-property drill-down."),
    ]
    for i, (title, body) in enumerate(features):
        _add_card(slide, 0.8 + i * 4.0, 2.5, 3.6, 2.5,
                  title, body, number=i + 1)

    _add_body(slide,
              "Ideal for acquisition teams evaluating multiple properties in a single transaction.",
              top=5.5, size=13)

    _add_footer(slide, 9)


def _slide_10_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "Report Comparison")
    _add_subtitle(slide, "Track changes over time with side-by-side analysis.")
    _add_rule(slide, 2.1)

    items = [
        ("Score Delta", "See how a property's activity score\nhas changed between reports."),
        ("Findings Diff", "Findings only in Report A, shared\nfindings, and findings only in Report B."),
        ("Record Changes", "Compare violation, inspection,\npermit, and lien counts over time."),
    ]
    for i, (title, body) in enumerate(items):
        _add_card(slide, 0.8 + i * 4.0, 2.5, 3.6, 2.5,
                  title, body, number=i + 1)

    _add_body(slide,
              "Monitor properties on your watchlist and catch emerging issues before closing.",
              top=5.5, size=13)

    _add_footer(slide, 10)


def _slide_11_why_civitas(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "Why CIVITAS")
    _add_subtitle(slide, "Purpose-built for transactional due diligence.")
    _add_rule(slide, 2.1)

    diffs = [
        ("Deterministic",
         "Every finding is computed by SQL rules\u2014never AI-guessed. "
         "Results are reproducible and auditable."),
        ("Transparent",
         "Every triggered finding includes a description, weight, "
         "and the supporting records that caused it."),
        ("Legally Cautious",
         "AI narratives avoid speculation, predictions, and legal advice. "
         "Clear disclaimers on every report."),
        ("Fast",
         "Full property analysis in under 60 seconds. Batch portfolios "
         "processed with real-time streaming progress."),
    ]
    for i, (title, body) in enumerate(diffs):
        col = i % 2
        row = i // 2
        _add_card(slide, 0.8 + col * 6.0, 2.4 + row * 2.1, 5.6, 1.7,
                  title, body, title_size=15, body_size=12)

    _add_footer(slide, 11)


def _slide_12_pricing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_title(slide, "Pricing")
    _add_subtitle(slide, "Simple, transparent plans that scale with your practice.")
    _add_rule(slide, 2.1)

    plans = [
        ("Starter", "$99/mo",
         "\u2022  10 reports per month\n"
         "\u2022  Full activity scoring & PDF export\n"
         "\u2022  AI narrative summaries\n"
         "\u2022  Email support"),
        ("Professional", "$249/mo",
         "\u2022  50 reports per month\n"
         "\u2022  Batch portfolio upload\n"
         "\u2022  Report comparison\n"
         "\u2022  Priority email support"),
        ("Enterprise", "Custom",
         "\u2022  Unlimited reports\n"
         "\u2022  API access & integrations\n"
         "\u2022  Dedicated account manager\n"
         "\u2022  Custom rule configuration"),
    ]
    for i, (name, price, features) in enumerate(plans):
        x = 0.8 + i * 4.0
        is_featured = (i == 1)

        # Card background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.4), Inches(3.6), Inches(3.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK if is_featured else CARD_BG
        shape.line.color.rgb = DARK_MID if is_featured else BORDER
        shape.line.width = Pt(1)
        shape.adjustments[0] = 0.04

        name_color = WHITE if is_featured else CHARCOAL
        price_color = WHITE if is_featured else CHARCOAL
        feat_color = TEXT_LIGHT if is_featured else TEXT_BODY

        # Plan name
        txBox = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(2.6), Inches(3), Inches(0.4)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name
        run.font.name = FONT
        run.font.size = Pt(16)
        run.font.color.rgb = name_color
        run.font.bold = True

        # Price
        txBox2 = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(3.0), Inches(3), Inches(0.5)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = price
        run2.font.name = FONT
        run2.font.size = Pt(28)
        run2.font.color.rgb = price_color
        run2.font.bold = True

        # Features
        txBox3 = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(3.7), Inches(3), Inches(2.3)
        )
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = features
        run3.font.name = FONT
        run3.font.size = Pt(12)
        run3.font.color.rgb = feat_color
        p3.line_spacing = Pt(22)

    # Popular badge on Professional (dark pill)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(2.2), Inches(1.4), Inches(0.28)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = CHARCOAL
    badge.line.fill.background()
    badge.adjustments[0] = 0.4
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "POPULAR"
    run.font.name = FONT
    run.font.size = Pt(9)
    run.font.color.rgb = WHITE
    run.font.bold = True

    _add_footer(slide, 12)


def _slide_13_cta(prs):
    """Get Started / CTA slide: dark hero."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)

    # Dark hero
    hero = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(4.5)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = DARK
    hero.line.fill.background()

    # CIVITAS wordmark
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.0), Inches(11.5), Inches(1)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "C I V I T A S"
    run.font.name = FONT
    run.font.size = Pt(48)
    run.font.color.rgb = WHITE
    run.font.bold = True

    # Subtitle
    txBox2 = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.5)
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Municipal Intelligence"
    run2.font.name = FONT
    run2.font.size = Pt(18)
    run2.font.color.rgb = TEXT_LIGHT

    # CTA
    txBox3 = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.6)
    )
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Request a Demo"
    run3.font.name = FONT
    run3.font.size = Pt(26)
    run3.font.color.rgb = WHITE
    run3.font.bold = True

    # Contact info
    txBox4 = slide.shapes.add_textbox(
        Inches(0.8), Inches(5.2), Inches(11.5), Inches(1)
    )
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = "contact@civitas.ai  \u2022  civitas.ai"
    run4.font.name = FONT
    run4.font.size = Pt(15)
    run4.font.color.rgb = TEXT_SEC

    p5 = tf4.add_paragraph()
    p5.alignment = PP_ALIGN.CENTER
    p5.space_before = Pt(8)
    run5 = p5.add_run()
    run5.text = "Chicago v1  \u2022  Deterministic  \u2022  Transparent  \u2022  Legally Cautious"
    run5.font.name = FONT
    run5.font.size = Pt(11)
    run5.font.color.rgb = TEXT_SEC

    _add_footer(slide, 13)


# -- Main ---------------------------------------------------------------------

def generate():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_01_title(prs)
    _slide_02_problem(prs)
    _slide_03_cost(prs)
    _slide_04_solution(prs)
    _slide_05_how_it_works(prs)
    _slide_06_activity_scoring(prs)
    _slide_07_report_deep_dive(prs)
    _slide_08_data_sources(prs)
    _slide_09_portfolio(prs)
    _slide_10_comparison(prs)
    _slide_11_why_civitas(prs)
    _slide_12_pricing(prs)
    _slide_13_cta(prs)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.dirname(script_dir)
    output_path = os.path.join(project_root, "civitas_pitch_deck.pptx")
    prs.save(output_path)
    print(f"Pitch deck saved to: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    generate()
